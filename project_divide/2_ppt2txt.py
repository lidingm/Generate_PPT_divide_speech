from pptx import Presentation
def ppt_to_word(ppt_file, word_file):
    # 读取PPT文件
    ppt = Presentation(ppt_file)

    with open(word_file, "w", encoding='UTF-8') as f:
        f.write('')

    with open(word_file, "a",encoding='UTF-8') as f:
        # 遍历PPT中的每一页
        for slide in ppt.slides:
            # 遍历幻灯片中的每一个形状
            for shape in slide.shapes:
                # 如果形状包含文本
                if hasattr(shape, "text"):
                    f.write(f'{shape.text}<!!>')
            f.write('\n///\n')

def generate_speech_prompt(previous):
    values = {'previous': previous,"填充自己的介绍、文章内容或文章背景的一个总体理解":"{填充自己的介绍、文章内容或文章背景的一个总体理解},",
              "进行填充演讲文稿大体":"{进行填充演讲文稿大体}"}
    with open("generate_speech_prompt.txt", "r", encoding='UTF-8') as f:
        data = f.read().format(**values)
        return data
if __name__ == '__main__':
    # 使用示例
    ppt_file = "../paper_test/AI.pptx"
    word_file = "ppt.txt"
    # 获取PPT每页中的文字，写入ppt.txt文件
    ppt_to_word(ppt_file, word_file)
    with open(word_file, "r", encoding='UTF-8') as f:
        data = f.read().replace('<!!>',' ')
    # 写入到4_speech_prompt.txt文件，其存放演讲文稿的prompt
    with open("4_speech_prompt.txt", "w", encoding='UTF-8') as f:
        f.write(generate_speech_prompt(data))